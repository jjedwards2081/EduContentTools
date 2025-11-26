"""
Game Manager module for handling Minecraft Education game folders and information.
"""

import os
import json
import shutil
import zipfile
import tempfile
import mimetypes
from datetime import datetime
import re


class GameManager:
    """Manages Minecraft Education game folders and information."""
    
    def __init__(self, settings):
        self.settings = settings
        # Store games in the same directory as the application
        app_dir = os.path.dirname(os.path.abspath(__file__))
        self.games_dir = os.path.join(app_dir, "games")
        self._ensure_games_dir()
    
    def _ensure_games_dir(self):
        """Ensure the games directory exists."""
        if not os.path.exists(self.games_dir):
            os.makedirs(self.games_dir)
    
    def create_game_folder(self, game_name):
        """Create a new game folder."""
        game_path = os.path.join(self.games_dir, game_name)
        
        if os.path.exists(game_path):
            return False
        
        try:
            os.makedirs(game_path)
            
            # Create initial metadata file
            metadata = {
                "name": game_name,
                "created": datetime.now().isoformat(),
                "modified": datetime.now().isoformat()
            }
            
            metadata_file = os.path.join(game_path, "metadata.json")
            with open(metadata_file, 'w') as f:
                json.dump(metadata, f, indent=4)
            
            return True
        except (OSError, IOError) as e:
            print(f"Error creating game folder: {e}")
            return False
    
    def delete_game_folder(self, game_name):
        """Delete a game folder and all its contents."""
        game_path = os.path.join(self.games_dir, game_name)
        
        if not os.path.exists(game_path):
            return False
        
        try:
            shutil.rmtree(game_path)
            return True
        except Exception as e:
            print(f"Error deleting game folder: {e}")
            return False
    
    def list_games(self):
        """List all game folders."""
        try:
            games = [d for d in os.listdir(self.games_dir) 
                    if os.path.isdir(os.path.join(self.games_dir, d))]
            return sorted(games)
        except OSError:
            return []
    
    def list_world_files_in_downloads(self):
        """List all .mcworld and .mctemplate files in the Downloads folder."""
        downloads_dir = os.path.join(os.path.expanduser("~"), "Downloads")
        
        if not os.path.exists(downloads_dir):
            return []
        
        try:
            world_files = []
            for file in os.listdir(downloads_dir):
                if file.endswith(('.mcworld', '.mctemplate')):
                    full_path = os.path.join(downloads_dir, file)
                    if os.path.isfile(full_path):
                        world_files.append(full_path)
            return sorted(world_files)
        except OSError:
            return []
    
    def upload_world_file(self, game_name, source_file):
        """Upload a world file to a game folder."""
        game_path = os.path.join(self.games_dir, game_name)
        
        if not os.path.exists(game_path):
            return False
        
        if not os.path.exists(source_file):
            return False
        
        try:
            # Get the filename
            filename = os.path.basename(source_file)
            
            # Create worlds subdirectory if it doesn't exist
            worlds_dir = os.path.join(game_path, "worlds")
            if not os.path.exists(worlds_dir):
                os.makedirs(worlds_dir)
            
            # Copy the file
            destination = os.path.join(worlds_dir, filename)
            shutil.copy2(source_file, destination)
            
            # Update metadata with world file info
            self._add_world_file_to_metadata(game_name, filename)
            
            return True
        except (OSError, IOError) as e:
            print(f"Error uploading world file: {e}")
            return False
    
    def _add_world_file_to_metadata(self, game_name, filename):
        """Add world file information to game metadata."""
        game_path = os.path.join(self.games_dir, game_name)
        metadata_file = os.path.join(game_path, "metadata.json")
        
        try:
            with open(metadata_file, 'r') as f:
                metadata = json.load(f)
            
            if "world_files" not in metadata:
                metadata["world_files"] = []
            
            # Add file info if not already present
            file_info = {
                "filename": filename,
                "uploaded": datetime.now().isoformat()
            }
            
            # Check if file already exists in metadata
            existing = [f for f in metadata["world_files"] if f["filename"] == filename]
            if not existing:
                metadata["world_files"].append(file_info)
            
            metadata["modified"] = datetime.now().isoformat()
            
            with open(metadata_file, 'w') as f:
                json.dump(metadata, f, indent=4)
        except (json.JSONDecodeError, IOError) as e:
            print(f"Error updating metadata: {e}")
    
    def save_game_info(self, game_name, info_type, content):
        """Save game information to a file."""
        game_path = os.path.join(self.games_dir, game_name)
        
        if not os.path.exists(game_path):
            return False
        
        info_file = os.path.join(game_path, f"{info_type}.json")
        
        try:
            data = {
                "type": info_type,
                "content": content,
                "modified": datetime.now().isoformat()
            }
            
            with open(info_file, 'w') as f:
                json.dump(data, f, indent=4)
            
            # Update metadata
            self._update_metadata(game_name)
            
            return True
        except (OSError, IOError) as e:
            print(f"Error saving game info: {e}")
            return False
    
    def load_game_info(self, game_name):
        """Load all game information."""
        game_path = os.path.join(self.games_dir, game_name)
        
        if not os.path.exists(game_path):
            return None
        
        info = {}
        info_types = ["context", "gameplay", "objectives", "lang_analysis", "document_analysis"]
        
        for info_type in info_types:
            info_file = os.path.join(game_path, f"{info_type}.json")
            if os.path.exists(info_file):
                try:
                    with open(info_file, 'r') as f:
                        data = json.load(f)
                        info[info_type] = data.get("content")
                except (json.JSONDecodeError, IOError):
                    pass
        
        # Load world files from metadata
        metadata_file = os.path.join(game_path, "metadata.json")
        if os.path.exists(metadata_file):
            try:
                with open(metadata_file, 'r') as f:
                    metadata = json.load(f)
                    if "world_files" in metadata:
                        info["world_files"] = metadata["world_files"]
            except (json.JSONDecodeError, IOError):
                pass
        
        return info if info else None
    
    def _update_metadata(self, game_name):
        """Update game metadata with current timestamp."""
        game_path = os.path.join(self.games_dir, game_name)
        metadata_file = os.path.join(game_path, "metadata.json")
        
        try:
            with open(metadata_file, 'r') as f:
                metadata = json.load(f)
            
            metadata["modified"] = datetime.now().isoformat()
            
            with open(metadata_file, 'w') as f:
                json.dump(metadata, f, indent=4)
        except (json.JSONDecodeError, IOError):
            pass
    
    def export_game(self, game_name):
        """Export game information to a markdown file."""
        info = self.load_game_info(game_name)
        
        if not info:
            return None
        
        export_dir = os.path.join(self.games_dir, game_name, "exports")
        if not os.path.exists(export_dir):
            os.makedirs(export_dir)
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        export_file = os.path.join(export_dir, f"{game_name}_{timestamp}.md")
        
        try:
            with open(export_file, 'w') as f:
                f.write(f"# {game_name}\n\n")
                f.write(f"*Exported: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n")
                
                if info:
                    if "world_files" in info and info["world_files"]:
                        f.write("## World Files\n\n")
                        for wf in info["world_files"]:
                            f.write(f"- **{wf['filename']}**\n")
                            f.write(f"  - Uploaded: {wf['uploaded']}\n")
                        f.write("\n")
                    
                    if "lang_analysis" in info and info["lang_analysis"]:
                        f.write("## Language File Analysis\n\n")
                        analysis = info["lang_analysis"]
                        if isinstance(analysis, dict):
                            if "analysis" in analysis:
                                f.write(f"{analysis['analysis']}\n\n")
                            elif "summary" in analysis:
                                f.write(f"{analysis['summary']}\n\n")
                        else:
                            f.write(f"{analysis}\n\n")
                    
                    if "document_analysis" in info and info["document_analysis"]:
                        f.write("## Document Analysis\n\n")
                        doc_analyses = info["document_analysis"]
                        if isinstance(doc_analyses, dict):
                            for doc_name, doc_data in doc_analyses.items():
                                f.write(f"### {doc_name}\n\n")
                                if isinstance(doc_data, dict) and "analysis" in doc_data:
                                    f.write(f"{doc_data['analysis']}\n\n")
                                elif isinstance(doc_data, dict) and "summary" in doc_data:
                                    f.write(f"{doc_data['summary']}\n\n")
                                else:
                                    f.write(f"{doc_data}\n\n")
                    
                    if "context" in info:
                        f.write("## Context\n\n")
                        f.write(f"{info['context']}\n\n")
                    
                    if "gameplay" in info:
                        f.write("## Gameplay Description\n\n")
                        f.write(f"{info['gameplay']}\n\n")
                    
                    if "objectives" in info:
                        f.write("## Learning Objectives\n\n")
                        if isinstance(info["objectives"], list):
                            for obj in info["objectives"]:
                                f.write(f"- {obj}\n")
                        else:
                            f.write(f"{info['objectives']}\n")
                        f.write("\n")
            
            return export_file
        except IOError as e:
            print(f"Error exporting game: {e}")
            return None
    
    def export_all_creations(self, game_name, export_format="md"):
        """Export all created resources to a single organized folder.
        
        Args:
            game_name: Name of the game
            export_format: Format for export - 'md', 'docx', or 'pdf'
        """
        game_path = os.path.join(self.games_dir, game_name)
        creations_dir = os.path.join(game_path, "creations")
        
        # Check if creations directory exists
        if not os.path.exists(creations_dir):
            return {
                "folder": None,
                "files": []
            }
        
        # Validate format
        export_format = export_format.lower()
        if export_format not in ['md', 'docx', 'pdf']:
            export_format = 'md'
        
        # Create export folder with timestamp
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        format_label = export_format.upper()
        export_folder_name = f"{game_name}_Creations_Export_{format_label}_{timestamp}"
        
        # Export to Downloads folder for easy access
        downloads_dir = os.path.join(os.path.expanduser("~"), "Downloads")
        export_path = os.path.join(downloads_dir, export_folder_name)
        
        try:
            # Create export directory
            os.makedirs(export_path)
            
            # Get all files from creations directory
            creation_files = []
            if os.path.exists(creations_dir):
                for filename in os.listdir(creations_dir):
                    file_path = os.path.join(creations_dir, filename)
                    if os.path.isfile(file_path) and filename.endswith('.md'):
                        # Convert and export based on format
                        if export_format == 'md':
                            # Copy markdown file as-is
                            shutil.copy2(file_path, os.path.join(export_path, filename))
                            creation_files.append(filename)
                        elif export_format == 'docx':
                            # Convert to Word
                            output_filename = filename.replace('.md', '.docx')
                            output_path = os.path.join(export_path, output_filename)
                            if self._convert_md_to_docx(file_path, output_path):
                                creation_files.append(output_filename)
                        elif export_format == 'pdf':
                            # Convert to PDF
                            output_filename = filename.replace('.md', '.pdf')
                            output_path = os.path.join(export_path, output_filename)
                            if self._convert_md_to_pdf(file_path, output_path):
                                creation_files.append(output_filename)
            
            # Create a README file in the export folder
            readme_path = os.path.join(export_path, "README.txt")
            with open(readme_path, 'w', encoding='utf-8') as f:
                f.write(f"Minecraft Education Content Export\n")
                f.write(f"{'=' * 50}\n\n")
                f.write(f"Game: {game_name}\n")
                f.write(f"Exported: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
                f.write(f"This folder contains all created educational resources:\n\n")
                
                # Categorize files
                student_files = [f for f in creation_files if f.startswith('Student_')]
                parent_files = [f for f in creation_files if f.startswith('Parent_')]
                teacher_files = [f for f in creation_files if f.startswith('Teacher_')]
                leadership_files = [f for f in creation_files if f.startswith('Leadership_')]
                curriculum_files = [f for f in creation_files if f.startswith('Curriculum_')]
                text_complexity_files = [f for f in creation_files if f.startswith('Text_Complexity_')]
                
                if student_files:
                    f.write("STUDENT RESOURCES:\n")
                    for file in student_files:
                        f.write(f"  - {file}\n")
                    f.write("\n")
                
                if parent_files:
                    f.write("PARENT RESOURCES:\n")
                    for file in parent_files:
                        f.write(f"  - {file}\n")
                    f.write("\n")
                
                if teacher_files:
                    f.write("TEACHER RESOURCES:\n")
                    for file in teacher_files:
                        f.write(f"  - {file}\n")
                    f.write("\n")
                
                if leadership_files:
                    f.write("LEADERSHIP RESOURCES:\n")
                    for file in leadership_files:
                        f.write(f"  - {file}\n")
                    f.write("\n")
                
                if curriculum_files:
                    f.write("CURRICULUM MAPPING:\n")
                    for file in curriculum_files:
                        f.write(f"  - {file}\n")
                    f.write("\n")
                
                if text_complexity_files:
                    f.write("TEXT COMPLEXITY ANALYSIS:\n")
                    for file in text_complexity_files:
                        f.write(f"  - {file}\n")
                    f.write("\n")
                
                f.write(f"\nTotal files: {len(creation_files)}\n")
                
                # Add format-specific information
                if export_format == 'md':
                    f.write(f"\nAll files are in Markdown (.md) format and can be opened with:\n")
                    f.write(f"  - Any text editor (Notepad, TextEdit, etc.)\n")
                    f.write(f"  - Markdown viewers (Typora, Marked, etc.)\n")
                    f.write(f"  - Word processors (Microsoft Word, Google Docs, etc.)\n")
                elif export_format == 'docx':
                    f.write(f"\nAll files are in Word Document (.docx) format and can be opened with:\n")
                    f.write(f"  - Microsoft Word\n")
                    f.write(f"  - Google Docs\n")
                    f.write(f"  - LibreOffice Writer\n")
                    f.write(f"  - Apple Pages\n")
                elif export_format == 'pdf':
                    f.write(f"\nAll files are in PDF (.pdf) format and can be opened with:\n")
                    f.write(f"  - Adobe Acrobat Reader\n")
                    f.write(f"  - Web browsers (Chrome, Safari, Firefox, etc.)\n")
                    f.write(f"  - Preview (macOS)\n")
                    f.write(f"  - Most PDF viewers\n")
            
            creation_files.append("README.txt")
            
            return {
                "folder": export_path,
                "files": creation_files
            }
        
        except Exception as e:
            print(f"Error exporting creations: {e}")
            return None
    
    def _convert_md_to_docx(self, md_file_path, output_path):
        """Convert a Markdown file to Word document format."""
        try:
            from docx import Document
            from docx.shared import Pt, RGBColor, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml.ns import qn
            from docx.oxml import OxmlElement
            
            # Read markdown content
            with open(md_file_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            # Create Word document
            doc = Document()
            
            # Process markdown line by line
            lines = md_content.split('\n')
            i = 0
            while i < len(lines):
                line = lines[i].rstrip()
                
                # Check for table (starts with |)
                if line.strip().startswith('|') and i + 1 < len(lines):
                    # Collect table rows
                    table_lines = []
                    while i < len(lines) and lines[i].strip().startswith('|'):
                        table_lines.append(lines[i].strip())
                        i += 1
                    
                    if len(table_lines) >= 2:
                        # Parse table
                        header_row = [cell.strip() for cell in table_lines[0].split('|')[1:-1]]
                        # Skip separator line
                        data_rows = []
                        for row_line in table_lines[2:]:
                            cells = [cell.strip() for cell in row_line.split('|')[1:-1]]
                            if cells:  # Only add non-empty rows
                                data_rows.append(cells)
                        
                        # Create Word table
                        if header_row:
                            table = doc.add_table(rows=1, cols=len(header_row))
                            table.style = 'Light Grid Accent 1'
                            
                            # Set table to auto-fit content and allow it to fit page
                            table.autofit = False
                            table.allow_autofit = True
                            
                            # Add header
                            hdr_cells = table.rows[0].cells
                            for idx, header in enumerate(header_row):
                                hdr_cells[idx].text = header
                                # Make header bold and set smaller font for better fit
                                for paragraph in hdr_cells[idx].paragraphs:
                                    paragraph.style = 'Normal'
                                    for run in paragraph.runs:
                                        run.font.bold = True
                                        run.font.size = Pt(10)
                                # Enable word wrap
                                hdr_cells[idx].width = Inches(6.5 / len(header_row))
                            
                            # Add data rows
                            for row_data in data_rows:
                                row_cells = table.add_row().cells
                                for idx, cell_data in enumerate(row_data):
                                    if idx < len(row_cells):
                                        row_cells[idx].text = cell_data
                                        # Set smaller font and enable wrapping
                                        for paragraph in row_cells[idx].paragraphs:
                                            for run in paragraph.runs:
                                                run.font.size = Pt(9)
                                        row_cells[idx].width = Inches(6.5 / len(header_row))
                            
                            doc.add_paragraph()  # Add space after table
                    continue
                
                # Headers
                if line.startswith('# '):
                    p = doc.add_heading(line[2:], level=1)
                elif line.startswith('## '):
                    p = doc.add_heading(line[3:], level=2)
                elif line.startswith('### '):
                    p = doc.add_heading(line[4:], level=3)
                elif line.startswith('#### '):
                    p = doc.add_heading(line[5:], level=4)
                # Bullet points
                elif line.startswith('- ') or line.startswith('* '):
                    doc.add_paragraph(line[2:], style='List Bullet')
                # Numbered lists
                elif re.match(r'^\d+\.\s', line):
                    text = re.sub(r'^\d+\.\s', '', line)
                    doc.add_paragraph(text, style='List Number')
                # Bold text (basic support)
                elif '**' in line:
                    p = doc.add_paragraph()
                    parts = line.split('**')
                    for j, part in enumerate(parts):
                        run = p.add_run(part)
                        if j % 2 == 1:  # Odd indices are bold
                            run.bold = True
                # Horizontal rule
                elif line.strip() == '---' or line.strip() == '***':
                    doc.add_paragraph('_' * 50)
                # Empty line
                elif not line.strip():
                    doc.add_paragraph()
                # Regular paragraph
                else:
                    doc.add_paragraph(line)
                
                i += 1
            
            # Save document
            doc.save(output_path)
            return True
            
        except ImportError:
            print(f"\n⚠️  python-docx package not installed. Install with: pip install python-docx")
            return False
        except Exception as e:
            print(f"Error converting to Word: {e}")
            return False
    
    def _format_bold_text(self, text):
        """Format bold markdown text (**text**) to reportlab bold tags (<b>text</b>)."""
        # Count occurrences of **
        count = text.count('**')
        if count == 0:
            return text
        
        # Replace in pairs
        result = text
        while '**' in result:
            # Replace first occurrence with <b>
            result = result.replace('**', '<b>', 1)
            # Replace second occurrence with </b>
            if '**' in result:
                result = result.replace('**', '</b>', 1)
            else:
                # Odd number of **, remove the unclosed tag
                result = result.replace('<b>', '**')
                break
        
        return result
    
    def _convert_md_to_pdf(self, md_file_path, output_path):
        """Convert a Markdown file to PDF format using reportlab."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
            from reportlab.lib.enums import TA_LEFT, TA_CENTER
            from reportlab.lib.colors import HexColor
            from reportlab.lib import colors
            
            # Read markdown content
            with open(md_file_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            # Create PDF document
            doc = SimpleDocTemplate(
                output_path,
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72
            )
            
            # Get styles
            styles = getSampleStyleSheet()
            
            # Create custom styles
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                textColor=HexColor('#2c3e50'),
                spaceAfter=20,
                spaceBefore=10
            )
            
            heading1_style = ParagraphStyle(
                'CustomHeading1',
                parent=styles['Heading1'],
                fontSize=18,
                textColor=HexColor('#34495e'),
                spaceAfter=12,
                spaceBefore=12
            )
            
            heading2_style = ParagraphStyle(
                'CustomHeading2',
                parent=styles['Heading2'],
                fontSize=14,
                textColor=HexColor('#555555'),
                spaceAfter=10,
                spaceBefore=10
            )
            
            body_style = ParagraphStyle(
                'CustomBody',
                parent=styles['BodyText'],
                fontSize=11,
                leading=16,
                spaceAfter=6
            )
            
            bullet_style = ParagraphStyle(
                'CustomBullet',
                parent=styles['BodyText'],
                fontSize=11,
                leading=16,
                leftIndent=20,
                spaceAfter=3
            )
            
            # Process markdown content
            story = []
            lines = md_content.split('\n')
            i = 0
            
            while i < len(lines):
                line = lines[i].rstrip()
                
                # Check for table (starts with |)
                if line.strip().startswith('|') and i + 1 < len(lines):
                    # Collect table rows
                    table_lines = []
                    while i < len(lines) and lines[i].strip().startswith('|'):
                        table_lines.append(lines[i].strip())
                        i += 1
                    
                    if len(table_lines) >= 2:
                        # Parse table
                        header_row = [cell.strip() for cell in table_lines[0].split('|')[1:-1]]
                        # Skip separator line (table_lines[1])
                        data_rows = []
                        for row_line in table_lines[2:]:
                            cells = [cell.strip() for cell in row_line.split('|')[1:-1]]
                            if cells:  # Only add non-empty rows
                                data_rows.append(cells)
                        
                        # Create PDF table
                        if header_row:
                            # Calculate available width (page width minus margins)
                            available_width = letter[0] - (2 * 72)  # 72 points = 1 inch margins
                            num_cols = len(header_row)
                            
                            # Parse raw table data first
                            raw_data = [header_row] + data_rows
                            
                            # Auto-calculate column widths based on content length
                            col_widths = []
                            for col_idx in range(num_cols):
                                max_len = max(
                                    len(str(raw_data[row_idx][col_idx])) 
                                    for row_idx in range(len(raw_data))
                                )
                                col_widths.append(max_len)
                            
                            # Normalize widths to fit available space
                            total_width = sum(col_widths)
                            if total_width > 0:
                                col_widths = [
                                    (w / total_width) * available_width 
                                    for w in col_widths
                                ]
                            else:
                                col_widths = [available_width / num_cols] * num_cols
                            
                            # Create styles for table cells
                            header_cell_style = ParagraphStyle(
                                'TableHeader',
                                parent=styles['Normal'],
                                fontSize=10,
                                leading=12,
                                textColor=colors.whitesmoke,
                                fontName='Helvetica-Bold',
                                alignment=TA_LEFT
                            )
                            
                            body_cell_style = ParagraphStyle(
                                'TableBody',
                                parent=styles['Normal'],
                                fontSize=9,
                                leading=11,
                                fontName='Helvetica',
                                alignment=TA_LEFT
                            )
                            
                            # Convert all cells to Paragraphs for proper word wrapping
                            table_data = []
                            
                            # Header row
                            header_paragraphs = [
                                Paragraph(str(cell), header_cell_style) 
                                for cell in header_row
                            ]
                            table_data.append(header_paragraphs)
                            
                            # Data rows
                            for row in data_rows:
                                row_paragraphs = [
                                    Paragraph(str(cell), body_cell_style) 
                                    for cell in row
                                ]
                                table_data.append(row_paragraphs)
                            
                            pdf_table = Table(table_data, colWidths=col_widths)
                            
                            # Style the table
                            pdf_table.setStyle(TableStyle([
                                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4472C4')),
                                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                                ('TOPPADDING', (0, 0), (-1, 0), 8),
                                ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                                ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#D9E2F3')),
                                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                                ('TOPPADDING', (0, 1), (-1, -1), 6),
                                ('BOTTOMPADDING', (0, 1), (-1, -1), 6),
                                ('LEFTPADDING', (0, 0), (-1, -1), 6),
                                ('RIGHTPADDING', (0, 0), (-1, -1), 6),
                                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.HexColor('#D9E2F3'), colors.white]),
                            ]))
                            
                            story.append(pdf_table)
                            story.append(Spacer(1, 0.2*inch))
                    continue
                
                # Skip empty lines
                if not line:
                    if story:  # Only add space if there's content
                        story.append(Spacer(1, 0.1*inch))
                    i += 1
                    continue
                
                # Headers
                if line.startswith('# '):
                    text = line[2:].strip()
                    story.append(Paragraph(text, title_style))
                elif line.startswith('## '):
                    text = line[3:].strip()
                    story.append(Paragraph(text, heading1_style))
                elif line.startswith('### '):
                    text = line[4:].strip()
                    story.append(Paragraph(text, heading2_style))
                elif line.startswith('#### '):
                    text = line[5:].strip()
                    story.append(Paragraph(text, heading2_style))
                
                # Bullet points
                elif line.startswith('- ') or line.startswith('* '):
                    text = '• ' + line[2:].strip()
                    text = self._format_bold_text(text)
                    story.append(Paragraph(text, bullet_style))
                
                # Numbered lists
                elif re.match(r'^\d+\.\s', line):
                    text = re.sub(r'^\d+\.\s', '', line)
                    text = self._format_bold_text(text)
                    story.append(Paragraph(text, bullet_style))
                
                # Horizontal rule
                elif line.strip() in ['---', '***', '___']:
                    story.append(Spacer(1, 0.2*inch))
                
                # Regular paragraph
                else:
                    text = self._format_bold_text(line)
                    story.append(Paragraph(text, body_style))
                
                i += 1
            
            # Build PDF
            doc.build(story)
            return True
            
        except ImportError:
            print(f"\n⚠️  reportlab package not installed. Install with: pip install reportlab")
            return False
        except Exception as e:
            print(f"Error converting to PDF: {e}")
            import traceback
            traceback.print_exc()
            return False
    
    def enhance_with_ai(self, text, enhancement_type):
        """Enhance text using Azure OpenAI."""
        if not self.settings.is_configured():
            print("Azure OpenAI is not configured!")
            return None
        
        try:
            from openai import AzureOpenAI
            
            config = self.settings.get_azure_config()
            
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            # Create appropriate prompt based on enhancement type
            prompts = {
                "context": "Enhance the following Minecraft Education game context. Make it more detailed, educational, and engaging while maintaining the core information:\n\n",
                "gameplay": "Improve the following gameplay description for a Minecraft Education game. Make it clearer, more detailed, and focused on educational aspects:\n\n",
                "objectives": "Refine these learning objectives for a Minecraft Education game. Make them specific, measurable, and aligned with educational standards:\n\n"
            }
            
            prompt = prompts.get(enhancement_type, "") + text
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "system", "content": "You are an educational content expert specializing in Minecraft Education."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=1000
            )
            
            return response.choices[0].message.content.strip()
        
        except Exception as e:
            print(f"Error enhancing with AI: {e}")
            return None
    
    def generate_context_from_data(self, game_name):
        """Generate game context from language file analysis and uploaded documents."""
        if not self.settings.is_configured():
            print("Azure OpenAI is not configured!")
            return None
        
        try:
            from openai import AzureOpenAI
            
            # Gather available data
            info = self.load_game_info(game_name)
            if not info:
                return None
            
            # Build prompt from available sources
            prompt_parts = []
            
            # Add existing gameplay if available (to maintain coherence)
            if "gameplay" in info and info["gameplay"]:
                prompt_parts.append(f"""EXISTING GAMEPLAY DESCRIPTION:
{info["gameplay"][:1500]}""")
            
            # Add language file analysis if available
            if "lang_analysis" in info and info["lang_analysis"]:
                lang_analysis = info["lang_analysis"]
                if isinstance(lang_analysis, dict) and "content" in lang_analysis:
                    lang_text = lang_analysis["content"]
                else:
                    lang_text = str(lang_analysis)
                
                prompt_parts.append(f"""LANGUAGE FILE ANALYSIS (NPC dialogue and in-game narrative):
{lang_text[:2000]}""")
            
            # Add document analysis if available
            if "document_analysis" in info and info["document_analysis"]:
                doc_analyses = info["document_analysis"]
                if isinstance(doc_analyses, dict):
                    for doc_name, doc_data in doc_analyses.items():
                        if isinstance(doc_data, dict) and "analysis" in doc_data:
                            analysis_text = doc_data["analysis"]
                        else:
                            analysis_text = str(doc_data)
                        
                        prompt_parts.append(f"""DOCUMENT: {doc_name}
{analysis_text[:2000]}""")
            
            if not prompt_parts:
                return None
            
            config = self.settings.get_azure_config()
            
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            full_prompt = f"""Based on the following information from a Minecraft Education game, generate a comprehensive game context that includes:

1. Educational Theme: The main educational focus and learning domain
2. Setting: The Minecraft world environment and scenario
3. Target Audience: Appropriate age/grade level
4. Duration: Estimated time to complete
5. Overview: A clear summary of what students will experience and learn

Available Information:

{chr(10).join(prompt_parts)}

IMPORTANT: If an existing gameplay description is provided above, ensure the game context aligns with and frames that gameplay appropriately. Generate a well-structured game context that synthesizes this information into a cohesive educational narrative."""
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "system", "content": "You are an expert in educational game design and Minecraft Education. Create comprehensive, well-organized game context descriptions that clearly communicate the educational value and gameplay experience."},
                    {"role": "user", "content": full_prompt}
                ],
                temperature=0.6,
                max_tokens=1500
            )
            
            return response.choices[0].message.content.strip()
        
        except Exception as e:
            print(f"Error generating context: {e}")
            return None
    
    def generate_gameplay_from_data(self, game_name):
        """Generate gameplay description from language file analysis and uploaded documents."""
        if not self.settings.is_configured():
            print("Azure OpenAI is not configured!")
            return None
        
        try:
            from openai import AzureOpenAI
            
            # Gather available data
            info = self.load_game_info(game_name)
            if not info:
                return None
            
            # Build prompt from available sources
            prompt_parts = []
            
            # Add existing context if available (to maintain coherence)
            if "context" in info and info["context"]:
                prompt_parts.append(f"""EXISTING GAME CONTEXT:
{info["context"][:1500]}""")
            
            # Add language file analysis if available
            if "lang_analysis" in info and info["lang_analysis"]:
                lang_analysis = info["lang_analysis"]
                if isinstance(lang_analysis, dict) and "content" in lang_analysis:
                    lang_text = lang_analysis["content"]
                else:
                    lang_text = str(lang_analysis)
                
                prompt_parts.append(f"""LANGUAGE FILE ANALYSIS (NPC dialogue and in-game narrative):
{lang_text[:2000]}""")
            
            # Add document analysis if available
            if "document_analysis" in info and info["document_analysis"]:
                doc_analyses = info["document_analysis"]
                if isinstance(doc_analyses, dict):
                    for doc_name, doc_data in doc_analyses.items():
                        if isinstance(doc_data, dict) and "analysis" in doc_data:
                            analysis_text = doc_data["analysis"]
                        else:
                            analysis_text = str(doc_data)
                        
                        prompt_parts.append(f"""DOCUMENT: {doc_name}
{analysis_text[:2000]}""")
            
            if not prompt_parts:
                return None
            
            config = self.settings.get_azure_config()
            
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            full_prompt = f"""Based on the following information from a Minecraft Education game, generate a comprehensive gameplay description that includes:

1. Setup: How players begin and what they're presented with
2. Core Mechanics: Main gameplay actions, interactions, and systems
3. Progression: How players advance through the experience
4. Challenges: Key obstacles, puzzles, or problems to solve
5. Success Criteria: How players know they've completed objectives

Available Information:

{chr(10).join(prompt_parts)}

IMPORTANT: If an existing game context is provided above, ensure the gameplay description aligns with and supports that educational context. Generate a detailed gameplay description that clearly explains what players will do, how they'll interact with the game world, and how the mechanics support the learning objectives."""
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "system", "content": "You are an expert in game design and educational technology, specializing in Minecraft Education. Create clear, detailed gameplay descriptions that explain mechanics, player actions, and learning integration."},
                    {"role": "user", "content": full_prompt}
                ],
                temperature=0.6,
                max_tokens=1500
            )
            
            return response.choices[0].message.content.strip()
        
        except Exception as e:
            print(f"Error generating gameplay description: {e}")
            return None
    
    def standardize_with_ai(self, text, content_type):
        """Standardize text using Azure OpenAI to ensure consistent format across all games."""
        if not self.settings.is_configured():
            print("Azure OpenAI is not configured!")
            return None
        
        try:
            from openai import AzureOpenAI
            
            config = self.settings.get_azure_config()
            
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            # Create standardization prompts for consistency
            prompts = {
                "context": """Standardize the following Minecraft Education game context into a consistent format with these sections:

1. Educational Theme: Brief description of the educational focus
2. Setting: The Minecraft world environment and setup
3. Target Audience: Age/grade level and prerequisite knowledge
4. Duration: Expected time to complete
5. Overview: A concise summary of the educational experience

Current context:
""",
                "gameplay": """Standardize the following gameplay description into a consistent format with these sections:

1. Setup: How players begin the game
2. Core Mechanics: Main gameplay actions and interactions
3. Progression: How players advance through the experience
4. Challenges: Key obstacles or problems to solve
5. Success Criteria: How players know they've completed the game

Current description:
""",
                "objectives": """Standardize the following learning objectives using this format:
- Start each objective with an action verb (understand, analyze, create, evaluate, etc.)
- Make them specific and measurable
- Align with Bloom's Taxonomy
- Keep each objective to one clear statement
- Format as a numbered list

Current objectives:
"""
            }
            
            prompt = prompts.get(content_type, "") + text
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "system", "content": "You are an educational content standardization expert. Your role is to format content consistently while preserving all important information. Maintain a professional, clear writing style suitable for educators."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3,  # Lower temperature for more consistent formatting
                max_tokens=1500
            )
            
            return response.choices[0].message.content.strip()
        
        except Exception as e:
            print(f"Error standardizing with AI: {e}")
            return None
    
    def test_azure_connection(self):
        """Test Azure OpenAI connection."""
        if not self.settings.is_configured():
            return False
        
        try:
            from openai import AzureOpenAI
            
            config = self.settings.get_azure_config()
            
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "user", "content": "Say 'Connection successful' if you can read this."}
                ],
                max_tokens=10
            )
            
            return True
        
        except Exception as e:
            print(f"Connection test failed: {e}")
            return False
    
    def extract_lang_files(self, game_name):
        """Extract language files (US/GB only) from uploaded .mcworld or .mctemplate files.
        Searches entire archive structure and intelligently selects the file most likely to contain NPC text."""
        game_path = os.path.join(self.games_dir, game_name)
        worlds_dir = os.path.join(game_path, "worlds")
        
        if not os.path.exists(worlds_dir):
            return {"success": False, "error": "No worlds directory found"}
        
        # Get all world files
        world_files = [f for f in os.listdir(worlds_dir) 
                      if f.endswith(('.mcworld', '.mctemplate'))]
        
        if not world_files:
            return {"success": False, "error": "No world files found"}
        
        # Create lang directory if it doesn't exist
        lang_dir = os.path.join(game_path, "lang")
        if not os.path.exists(lang_dir):
            os.makedirs(lang_dir)
        
        # Try to extract lang files from the first world file
        world_file_path = os.path.join(worlds_dir, world_files[0])
        
        try:
            with zipfile.ZipFile(world_file_path, 'r') as zip_ref:
                # Search for ALL .lang files in the entire archive structure
                all_files = zip_ref.namelist()
                lang_files = [f for f in all_files if f.endswith('.lang')]
                
                if not lang_files:
                    return {"success": False, "error": "No .lang files found in world archive"}
                
                # Filter for en_US and en_GB only
                preferred_langs = ['en_US.lang', 'en_GB.lang']
                us_gb_files = []
                
                for lang_file in lang_files:
                    for pref_lang in preferred_langs:
                        if lang_file.endswith(pref_lang):
                            us_gb_files.append(lang_file)
                            break
                
                if not us_gb_files:
                    return {"success": False, "error": "No en_US or en_GB language files found"}
                
                # Analyze each lang file to find the one with most NPC/dialogue content
                best_file = None
                best_score = 0
                best_lang = None
                all_candidates = []
                
                for lang_file_path in us_gb_files:
                    try:
                        lang_content = zip_ref.read(lang_file_path).decode('utf-8', errors='ignore')
                        lang_data = {}
                        
                        for line in lang_content.split('\n'):
                            line = line.strip()
                            if line and not line.startswith('#') and '=' in line:
                                key, value = line.split('=', 1)
                                lang_data[key.strip()] = value.strip()
                        
                        # Score this file based on likelihood of containing NPC text
                        score = self._score_lang_file_for_npc_content(lang_data, lang_file_path)
                        
                        all_candidates.append({
                            'path': lang_file_path,
                            'score': score,
                            'entries': len(lang_data),
                            'data': lang_data
                        })
                        
                        if score > best_score:
                            best_score = score
                            best_file = lang_file_path
                            best_lang = lang_data
                    
                    except Exception as e:
                        print(f"Error reading {lang_file_path}: {e}")
                        continue
                
                if not best_file:
                    return {"success": False, "error": "Could not parse any language files"}
                
                # Extract the best lang file
                lang_filename = best_file.split('/')[-1]
                output_file = os.path.join(lang_dir, lang_filename)
                
                with open(output_file, 'w', encoding='utf-8') as f:
                    for key, value in best_lang.items():
                        f.write(f"{key}={value}\n")
                
                # Save parsed data as JSON for easy access
                json_file = os.path.join(lang_dir, f"{lang_filename.replace('.lang', '.json')}")
                with open(json_file, 'w', encoding='utf-8') as f:
                    json.dump(best_lang, f, indent=4, ensure_ascii=False)
                
                # Save analysis info
                analysis_file = os.path.join(lang_dir, "extraction_analysis.json")
                with open(analysis_file, 'w', encoding='utf-8') as f:
                    json.dump({
                        'selected_file': best_file,
                        'score': best_score,
                        'candidates': [{
                            'path': c['path'],
                            'score': c['score'],
                            'entries': c['entries']
                        } for c in all_candidates]
                    }, f, indent=4)
                
                # Update metadata
                self._add_lang_file_to_metadata(game_name, lang_filename, len(best_lang))
                
                return {
                    "success": True,
                    "language": lang_filename.replace('.lang', ''),
                    "lang_file": output_file,
                    "entry_count": len(best_lang),
                    "preview": best_lang,
                    "selected_from": len(all_candidates),
                    "score": best_score,
                    "file_path": best_file
                }
        
        except zipfile.BadZipFile:
            return {"success": False, "error": "Invalid or corrupted world file"}
        except Exception as e:
            return {"success": False, "error": f"Error extracting lang files: {str(e)}"}
    
    def _score_lang_file_for_npc_content(self, lang_data, file_path):
        """Score a language file based on likelihood of containing NPC dialogue and educational content.
        Higher scores indicate more relevant content."""
        score = 0
        
        # Keywords that indicate NPC dialogue, educational content, or custom content
        high_value_keywords = [
            'npc', 'dialogue', 'dialog', 'quest', 'objective', 'mission',
            'instruction', 'tutorial', 'hint', 'guide', 'student', 'teacher',
            'learn', 'lesson', 'activity', 'challenge', 'task', 'step'
        ]
        
        medium_value_keywords = [
            'talk', 'speak', 'say', 'tell', 'ask', 'answer', 'explain',
            'description', 'info', 'message', 'text', 'story', 'narrative'
        ]
        
        # File path indicators (behavior packs, custom content)
        if 'behavior_pack' in file_path.lower() or 'bp' in file_path.lower():
            score += 50
        if 'resource_pack' in file_path.lower() or 'rp' in file_path.lower():
            score += 30
        if 'texts' in file_path.lower():
            score += 20
        
        # Analyze content
        for key, value in lang_data.items():
            key_lower = key.lower()
            value_lower = value.lower()
            
            # High value indicators
            for keyword in high_value_keywords:
                if keyword in key_lower or keyword in value_lower:
                    score += 5
            
            # Medium value indicators
            for keyword in medium_value_keywords:
                if keyword in key_lower or keyword in value_lower:
                    score += 2
            
            # Long text values likely indicate dialogue or descriptions
            if len(value) > 50:
                score += 3
            elif len(value) > 100:
                score += 5
            
            # Custom keys (not standard minecraft keys)
            if not key_lower.startswith(('entity.', 'item.', 'tile.', 'block.')):
                score += 2
            
            # Sentence-like content (has punctuation)
            if any(p in value for p in ['.', '?', '!']):
                score += 2
        
        # Bonus for larger files with diverse content
        entry_count = len(lang_data)
        if entry_count > 100:
            score += 20
        elif entry_count > 50:
            score += 10
        elif entry_count > 20:
            score += 5
        
        return score
    
    def _add_lang_file_to_metadata(self, game_name, lang_file, entry_count):
        """Add language file information to game metadata."""
        game_path = os.path.join(self.games_dir, game_name)
        metadata_file = os.path.join(game_path, "metadata.json")
        
        try:
            with open(metadata_file, 'r') as f:
                metadata = json.load(f)
            
            metadata["lang_file"] = {
                "filename": lang_file,
                "entry_count": entry_count,
                "extracted": datetime.now().isoformat()
            }
            
            metadata["modified"] = datetime.now().isoformat()
            
            with open(metadata_file, 'w') as f:
                json.dump(metadata, f, indent=4)
        except (json.JSONDecodeError, IOError) as e:
            print(f"Error updating metadata: {e}")
    
    def _load_metadata(self, game_name):
        """Load game metadata."""
        game_path = os.path.join(self.games_dir, game_name)
        metadata_file = os.path.join(game_path, "metadata.json")
        
        if os.path.exists(metadata_file):
            try:
                with open(metadata_file, 'r') as f:
                    return json.load(f)
            except (json.JSONDecodeError, IOError):
                return None
        return None
    
    def list_documents_in_downloads(self):
        """List all PDF, Word, and PowerPoint files in the Downloads folder."""
        downloads_dir = os.path.join(os.path.expanduser("~"), "Downloads")
        
        if not os.path.exists(downloads_dir):
            return []
        
        try:
            documents = []
            supported_extensions = ('.pdf', '.docx', '.doc', '.pptx', '.ppt')
            
            for file in os.listdir(downloads_dir):
                if file.lower().endswith(supported_extensions):
                    full_path = os.path.join(downloads_dir, file)
                    if os.path.isfile(full_path):
                        documents.append(full_path)
            return sorted(documents)
        except OSError:
            return []
    
    def upload_document(self, game_name, source_file):
        """Upload a document (PDF, Word, PPT) to a game folder."""
        game_path = os.path.join(self.games_dir, game_name)
        
        if not os.path.exists(game_path):
            return False
        
        if not os.path.exists(source_file):
            return False
        
        try:
            # Get the filename and extension
            filename = os.path.basename(source_file)
            ext = os.path.splitext(filename)[1].lower()
            
            # Create documents subdirectory if it doesn't exist
            docs_dir = os.path.join(game_path, "documents")
            if not os.path.exists(docs_dir):
                os.makedirs(docs_dir)
            
            # Copy the file
            destination = os.path.join(docs_dir, filename)
            shutil.copy2(source_file, destination)
            
            # Determine document type
            doc_type = "PDF" if ext == ".pdf" else "Word" if ext in [".docx", ".doc"] else "PowerPoint"
            
            # Update metadata with document info
            self._add_document_to_metadata(game_name, filename, doc_type)
            
            return True
        except (OSError, IOError) as e:
            print(f"Error uploading document: {e}")
            return False
    
    def _add_document_to_metadata(self, game_name, filename, doc_type):
        """Add document information to game metadata."""
        game_path = os.path.join(self.games_dir, game_name)
        metadata_file = os.path.join(game_path, "metadata.json")
        
        try:
            with open(metadata_file, 'r') as f:
                metadata = json.load(f)
            
            if "documents" not in metadata:
                metadata["documents"] = []
            
            # Add document info if not already present
            doc_info = {
                "filename": filename,
                "type": doc_type,
                "uploaded": datetime.now().isoformat(),
                "ai_analyzed": False
            }
            
            # Check if file already exists in metadata
            existing = [d for d in metadata["documents"] if d["filename"] == filename]
            if not existing:
                metadata["documents"].append(doc_info)
            
            metadata["modified"] = datetime.now().isoformat()
            
            with open(metadata_file, 'w') as f:
                json.dump(metadata, f, indent=4)
        except (json.JSONDecodeError, IOError) as e:
            print(f"Error updating metadata: {e}")
    
    def regenerate_document_analysis(self, game_name):
        """Regenerate document analysis after a document is removed.
        Re-analyzes remaining documents to update the combined analysis."""
        if not self.settings.is_configured():
            return False
        
        game_path = os.path.join(self.games_dir, game_name)
        metadata_file = os.path.join(game_path, "metadata.json")
        
        try:
            # Load metadata to get remaining documents
            if os.path.exists(metadata_file):
                with open(metadata_file, 'r') as f:
                    metadata = json.load(f)
                
                if "documents" in metadata and metadata["documents"]:
                    remaining_docs = metadata["documents"]
                    
                    print(f"\nUpdating analysis for {len(remaining_docs)} remaining document(s)...")
                    
                    # Re-analyze each remaining document that was previously analyzed
                    for doc in remaining_docs:
                        if doc.get("ai_analyzed"):
                            print(f"   Re-analyzing {doc['filename']}...")
                            self.analyze_document_with_ai(game_name, doc['filename'])
                    
                    print("✓ Analysis updated!")
                    return True
            
            return True
        except Exception as e:
            print(f"Error regenerating analysis: {e}")
            return False
    
    def analyze_document_with_ai(self, game_name, filename):
        """Analyze uploaded document using Azure OpenAI and save in standardized format."""
        if not self.settings.is_configured():
            print("Azure OpenAI is not configured!")
            return None
        
        game_path = os.path.join(self.games_dir, game_name)
        docs_dir = os.path.join(game_path, "documents")
        doc_path = os.path.join(docs_dir, filename)
        
        if not os.path.exists(doc_path):
            return None
        
        try:
            # Extract text from document
            text_content = self._extract_text_from_document(doc_path)
            
            if not text_content:
                print("Could not extract text from document.")
                return None
            
            # Limit text size for API call (take first ~4000 words)
            words = text_content.split()
            if len(words) > 4000:
                text_content = ' '.join(words[:4000]) + "\\n\\n[Document truncated for analysis]"
            
            from openai import AzureOpenAI
            
            config = self.settings.get_azure_config()
            
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            prompt = f"""Analyze this educational document related to a Minecraft Education game and provide a structured analysis:

1. DOCUMENT OVERVIEW:
   - Main purpose and audience
   - Document type and structure
   
2. EDUCATIONAL CONTENT:
   - Key learning objectives mentioned
   - Subject areas and concepts covered
   - Grade level and curriculum connections
   
3. GAME-RELATED INFORMATION:
   - Game mechanics or features described
   - Gameplay instructions or scenarios
   - Educational activities or challenges
   
4. TEACHER/FACILITATOR GUIDANCE:
   - Implementation instructions
   - Assessment criteria or rubrics
   - Time requirements and prerequisites
   
5. ADDITIONAL RESOURCES:
   - Related materials or extensions
   - Standards alignment
   - Supplementary activities

Document content:
{text_content}"""
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "system", "content": "You are an educational content analyst specializing in Minecraft Education. Provide detailed, structured analysis of educational documents to extract teaching guidance, learning objectives, and implementation details."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.5,
                max_tokens=2500
            )
            
            analysis_text = response.choices[0].message.content.strip()
            
            # Load existing document analysis or create new
            analysis_file = os.path.join(game_path, "document_analysis.json")
            if os.path.exists(analysis_file):
                with open(analysis_file, 'r') as f:
                    all_analyses = json.load(f)
                    if "content" in all_analyses:
                        all_analyses = all_analyses["content"]
            else:
                all_analyses = {}
            
            # Save analysis for this document
            all_analyses[filename] = {
                "analyzed_at": datetime.now().isoformat(),
                "document_type": os.path.splitext(filename)[1],
                "analysis": analysis_text,
                "summary": self._extract_summary(analysis_text)
            }
            
            # Save to game info folder
            self.save_game_info(game_name, "document_analysis", all_analyses)
            
            # Update metadata
            self._update_document_metadata(game_name, filename, analyzed=True)
            
            return analysis_text
        
        except Exception as e:
            print(f"Error analyzing document: {e}")
            return None
    
    def _extract_text_from_document(self, file_path):
        """Extract text content from PDF, Word, or PowerPoint files."""
        ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if ext == '.pdf':
                return self._extract_text_from_pdf(file_path)
            elif ext in ['.docx', '.doc']:
                return self._extract_text_from_word(file_path)
            elif ext in ['.pptx', '.ppt']:
                return self._extract_text_from_ppt(file_path)
        except Exception as e:
            print(f"Error extracting text: {e}")
            return None
    
    def _extract_text_from_pdf(self, file_path):
        """Extract text from PDF file."""
        try:
            import PyPDF2
            text = []
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text.append(page.extract_text())
            return '\\n'.join(text)
        except ImportError:
            print("PyPDF2 not installed. Install with: pip install PyPDF2")
            return None
        except Exception as e:
            print(f"Error reading PDF: {e}")
            return None
    
    def _extract_text_from_word(self, file_path):
        """Extract text from Word document."""
        try:
            from docx import Document
            doc = Document(file_path)
            text = []
            for paragraph in doc.paragraphs:
                text.append(paragraph.text)
            return '\\n'.join(text)
        except ImportError:
            print("python-docx not installed. Install with: pip install python-docx")
            return None
        except Exception as e:
            print(f"Error reading Word document: {e}")
            return None
    
    def _extract_text_from_ppt(self, file_path):
        """Extract text from PowerPoint presentation."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\\n'.join(text)
        except ImportError:
            print("python-pptx not installed. Install with: pip install python-pptx")
            return None
        except Exception as e:
            print(f"Error reading PowerPoint: {e}")
            return None
    
    def _update_document_metadata(self, game_name, filename, analyzed=False):
        """Update document metadata with analysis status."""
        game_path = os.path.join(self.games_dir, game_name)
        metadata_file = os.path.join(game_path, "metadata.json")
        
        try:
            with open(metadata_file, 'r') as f:
                metadata = json.load(f)
            
            if "documents" in metadata:
                for doc in metadata["documents"]:
                    if doc["filename"] == filename:
                        doc["ai_analyzed"] = analyzed
                        doc["analyzed_at"] = datetime.now().isoformat()
                        break
            
            metadata["modified"] = datetime.now().isoformat()
            
            with open(metadata_file, 'w') as f:
                json.dump(metadata, f, indent=4)
        except (json.JSONDecodeError, IOError) as e:
            print(f"Error updating metadata: {e}")
    
    def remove_document(self, game_name, filename):
        """Remove a specific document and its analysis."""
        game_path = os.path.join(self.games_dir, game_name)
        
        try:
            # Check if document was analyzed before removal
            metadata_file = os.path.join(game_path, "metadata.json")
            was_analyzed = False
            has_remaining_docs = False
            
            if os.path.exists(metadata_file):
                with open(metadata_file, 'r') as f:
                    metadata = json.load(f)
                
                if "documents" in metadata:
                    for doc in metadata["documents"]:
                        if doc["filename"] == filename:
                            was_analyzed = doc.get("ai_analyzed", False)
                    
                    # Count remaining documents after removal
                    remaining = [d for d in metadata["documents"] if d["filename"] != filename]
                    has_remaining_docs = len(remaining) > 0
            
            # Remove physical file
            docs_dir = os.path.join(game_path, "documents")
            doc_file = os.path.join(docs_dir, filename)
            if os.path.exists(doc_file):
                os.remove(doc_file)
            
            # Remove from metadata
            if os.path.exists(metadata_file):
                with open(metadata_file, 'r') as f:
                    metadata = json.load(f)
                
                if "documents" in metadata:
                    metadata["documents"] = [d for d in metadata["documents"] if d["filename"] != filename]
                    metadata["modified"] = datetime.now().isoformat()
                    
                    with open(metadata_file, 'w') as f:
                        json.dump(metadata, f, indent=4)
            
            # Remove from document_analysis
            analysis_file = os.path.join(game_path, "document_analysis.json")
            if os.path.exists(analysis_file):
                with open(analysis_file, 'r') as f:
                    analysis_data = json.load(f)
                
                if "content" in analysis_data:
                    content = analysis_data["content"]
                    if filename in content:
                        del content[filename]
                        
                        # Save updated analysis
                        if content:  # If there are still other analyses
                            self.save_game_info(game_name, "document_analysis", content)
                        else:  # If this was the last document
                            os.remove(analysis_file)
            
            # If document was analyzed and there are remaining documents, regenerate analysis
            if was_analyzed and has_remaining_docs:
                self.regenerate_document_analysis(game_name)
            
            return True
        except Exception as e:
            print(f"Error removing document: {e}")
            return False
    
    def remove_all_documents(self, game_name):
        """Remove all documents and their analysis."""
        game_path = os.path.join(self.games_dir, game_name)
        
        try:
            # Remove all physical files
            docs_dir = os.path.join(game_path, "documents")
            if os.path.exists(docs_dir):
                shutil.rmtree(docs_dir)
            
            # Remove from metadata
            metadata_file = os.path.join(game_path, "metadata.json")
            if os.path.exists(metadata_file):
                with open(metadata_file, 'r') as f:
                    metadata = json.load(f)
                
                if "documents" in metadata:
                    metadata["documents"] = []
                    metadata["modified"] = datetime.now().isoformat()
                    
                    with open(metadata_file, 'w') as f:
                        json.dump(metadata, f, indent=4)
            
            # Remove document_analysis file (no regeneration needed since all removed)
            analysis_file = os.path.join(game_path, "document_analysis.json")
            if os.path.exists(analysis_file):
                os.remove(analysis_file)
            
            return True
        except Exception as e:
            print(f"Error removing all documents: {e}")
            return False
    
    def _gather_all_game_info(self, game_name):
        """Gather all available information about a game for creation tasks."""
        info = self.load_game_info(game_name)
        metadata = self._load_metadata(game_name)
        
        combined_info = {
            "game_name": game_name,
            "context": info.get("context", "") if info else "",
            "gameplay": info.get("gameplay", "") if info else "",
            "objectives": info.get("objectives", []) if info else [],
            "lang_analysis": "",
            "document_analysis": "",
            "has_world_files": False
        }
        
        # Add language analysis
        if info and "lang_analysis" in info:
            lang_data = info["lang_analysis"]
            if isinstance(lang_data, dict) and "analysis" in lang_data:
                combined_info["lang_analysis"] = lang_data["analysis"]
        
        # Add document analysis
        if info and "document_analysis" in info:
            doc_data = info["document_analysis"]
            analyses = []
            if isinstance(doc_data, dict):
                for doc_name, doc_analysis in doc_data.items():
                    if isinstance(doc_analysis, dict) and "analysis" in doc_analysis:
                        analyses.append(f"From {doc_name}:\n{doc_analysis['analysis']}")
            if analyses:
                combined_info["document_analysis"] = "\n\n".join(analyses)
        
        # Check for world files
        if info and "world_files" in info and info["world_files"]:
            combined_info["has_world_files"] = True
        
        return combined_info
    
    def create_student_guide(self, game_name):
        """Create a comprehensive student guide using all available information."""
        if not self.settings.is_configured():
            return None
        
        game_info = self._gather_all_game_info(game_name)
        
        if not game_info["context"] and not game_info["gameplay"]:
            print("Not enough game information available. Please add context and gameplay first.")
            return None
        
        try:
            from openai import AzureOpenAI
            
            config = self.settings.get_azure_config()
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            # Build comprehensive prompt
            objectives_text = "\n".join([f"- {obj}" for obj in game_info["objectives"]]) if game_info["objectives"] else "Not specified"
            
            prompt = f"""Create a comprehensive STUDENT GUIDE for the Minecraft Education game "{game_info['game_name']}". 

This guide should be written directly for students and should be clear, engaging, and easy to follow.

Available Information:

GAME CONTEXT:
{game_info['context']}

GAMEPLAY:
{game_info['gameplay']}

LEARNING OBJECTIVES:
{objectives_text}

LANGUAGE FILE ANALYSIS:
{game_info['lang_analysis'] if game_info['lang_analysis'] else 'Not available'}

ADDITIONAL CONTEXT:
{game_info['document_analysis'] if game_info['document_analysis'] else 'Not available'}

Create a student guide with the following sections:

1. WELCOME & INTRODUCTION
   - Engaging introduction to the game
   - What students will learn and do

2. GETTING STARTED
   - How to begin the game
   - Basic controls and navigation
   - What to expect

3. GAME WALKTHROUGH
   - Step-by-step guidance through the game
   - Key objectives and milestones
   - Tips and strategies

4. LEARNING POINTS
   - Important concepts to understand
   - How the game relates to real-world learning

5. CHALLENGES & SOLUTIONS
   - Common challenges students might face
   - How to overcome them

6. REFLECTION QUESTIONS
   - Questions to think about while playing
   - Discussion prompts

7. ADDITIONAL RESOURCES
   - Where to learn more
   - Related topics to explore

Make it student-friendly, encouraging, and educational!"""
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "system", "content": "You are an expert educational content creator specializing in creating engaging student materials for Minecraft Education. Write in a friendly, encouraging tone that speaks directly to students."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=3000
            )
            
            guide_content = response.choices[0].message.content.strip()
            
            # Save to creations directory
            game_path = os.path.join(self.games_dir, game_name)
            creations_dir = os.path.join(game_path, "creations")
            if not os.path.exists(creations_dir):
                os.makedirs(creations_dir)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            guide_file = os.path.join(creations_dir, f"Student_Guide_{timestamp}.md")
            
            with open(guide_file, 'w', encoding='utf-8') as f:
                f.write(f"# Student Guide: {game_info['game_name']}\n\n")
                f.write(f"*Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n")
                f.write("---\n\n")
                f.write(guide_content)
            
            return guide_file
        
        except Exception as e:
            print(f"Error creating student guide: {e}")
            return None
    
    def create_student_workbook(self, game_name):
        """Create an interactive student workbook with activities."""
        if not self.settings.is_configured():
            return None
        
        game_info = self._gather_all_game_info(game_name)
        
        if not game_info["context"] and not game_info["gameplay"]:
            print("Not enough game information available. Please add context and gameplay first.")
            return None
        
        try:
            from openai import AzureOpenAI
            
            config = self.settings.get_azure_config()
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            objectives_text = "\n".join([f"- {obj}" for obj in game_info["objectives"]]) if game_info["objectives"] else "Not specified"
            
            prompt = f"""Create an interactive STUDENT WORKBOOK for the Minecraft Education game "{game_info['game_name']}". 

This workbook should include spaces for students to write, draw, and reflect. Use clear formatting with sections for student responses.

Available Information:

GAME CONTEXT:
{game_info['context']}

GAMEPLAY:
{game_info['gameplay']}

LEARNING OBJECTIVES:
{objectives_text}

LANGUAGE FILE ANALYSIS:
{game_info['lang_analysis'] if game_info['lang_analysis'] else 'Not available'}

ADDITIONAL CONTEXT:
{game_info['document_analysis'] if game_info['document_analysis'] else 'Not available'}

Create a workbook with these sections:

1. MY GAME JOURNAL
   - Space for students to record their progress
   - What I learned today sections

2. PRE-GAME ACTIVITIES
   - Background knowledge activation
   - Prediction exercises
   - Vocabulary preview

3. DURING GAMEPLAY
   - Observation notes
   - Screenshot analysis prompts
   - Strategy planning sections
   - Problem-solving worksheets

4. CONCEPT EXPLORATION
   - Activities related to learning objectives
   - Draw, label, and explain exercises
   - Real-world connections

5. CHALLENGES & ACHIEVEMENTS
   - Track completed objectives
   - Describe challenges faced
   - Solutions tried

6. POST-GAME REFLECTION
   - What I learned
   - How it connects to class
   - Questions I still have

7. EXTENSION ACTIVITIES
   - Creative projects
   - Research prompts
   - Design challenges

Use formatting like:
- [Write your answer here]
- _____________________________ (for fill-in-the-blank)
- [ ] Checkboxes for tracking
- Numbered/lettered spaces for responses"""
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "system", "content": "You are an expert educational content creator. Create engaging, interactive workbooks that encourage active learning and reflection. Include clear spaces for student responses."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=3500
            )
            
            workbook_content = response.choices[0].message.content.strip()
            
            # Save to creations directory
            game_path = os.path.join(self.games_dir, game_name)
            creations_dir = os.path.join(game_path, "creations")
            if not os.path.exists(creations_dir):
                os.makedirs(creations_dir)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            workbook_file = os.path.join(creations_dir, f"Student_Workbook_{timestamp}.md")
            
            with open(workbook_file, 'w', encoding='utf-8') as f:
                f.write(f"# Student Workbook: {game_info['game_name']}\n\n")
                f.write(f"*Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n")
                f.write(f"**Student Name:** _____________________________\n\n")
                f.write(f"**Date:** _____________________________\n\n")
                f.write("---\n\n")
                f.write(workbook_content)
            
            return workbook_file
        
        except Exception as e:
            print(f"Error creating student workbook: {e}")
            return None
    
    def create_student_quiz(self, game_name):
        """Create a student quiz with answer key."""
        if not self.settings.is_configured():
            return None
        
        game_info = self._gather_all_game_info(game_name)
        
        if not game_info["context"] and not game_info["gameplay"]:
            print("Not enough game information available. Please add context and gameplay first.")
            return None
        
        try:
            from openai import AzureOpenAI
            
            config = self.settings.get_azure_config()
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            objectives_text = "\n".join([f"- {obj}" for obj in game_info["objectives"]]) if game_info["objectives"] else "Not specified"
            
            prompt = f"""Create a comprehensive STUDENT QUIZ for the Minecraft Education game "{game_info['game_name']}". 

Available Information:

GAME CONTEXT:
{game_info['context']}

GAMEPLAY:
{game_info['gameplay']}

LEARNING OBJECTIVES:
{objectives_text}

LANGUAGE FILE ANALYSIS:
{game_info['lang_analysis'] if game_info['lang_analysis'] else 'Not available'}

ADDITIONAL CONTEXT:
{game_info['document_analysis'] if game_info['document_analysis'] else 'Not available'}

Create a quiz with:

PART 1: MULTIPLE CHOICE (10 questions)
- Cover key concepts from the game
- 4 options each (A, B, C, D)
- Mix of difficulty levels

PART 2: TRUE/FALSE (5 questions)
- Test understanding of game mechanics and concepts

PART 3: SHORT ANSWER (5 questions)
- Require brief written responses
- Test application of knowledge

PART 4: REFLECTION (2 questions)
- Open-ended questions about learning
- Connect game to real-world concepts

Format the quiz clearly with numbered questions. Then create a SEPARATE ANSWER KEY.

Provide both:
1. The student quiz (questions only)
2. Complete answer key with explanations"""
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "system", "content": "You are an expert educational assessment creator. Create clear, fair quizzes that assess student understanding at multiple levels. Provide detailed answer keys with explanations."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.6,
                max_tokens=3500
            )
            
            full_content = response.choices[0].message.content.strip()
            
            # Try to split into quiz and answers
            # Look for common separators
            separators = ["ANSWER KEY", "Answer Key", "ANSWERS", "Answers:", "---ANSWERS---"]
            quiz_content = full_content
            answer_content = ""
            
            for separator in separators:
                if separator in full_content:
                    parts = full_content.split(separator, 1)
                    quiz_content = parts[0].strip()
                    answer_content = separator + parts[1].strip()
                    break
            
            # If no separator found, create a simple split
            if not answer_content:
                quiz_content = full_content
                answer_content = "Answer key included in quiz document."
            
            # Save to creations directory
            game_path = os.path.join(self.games_dir, game_name)
            creations_dir = os.path.join(game_path, "creations")
            if not os.path.exists(creations_dir):
                os.makedirs(creations_dir)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            quiz_file = os.path.join(creations_dir, f"Student_Quiz_{timestamp}.md")
            answers_file = os.path.join(creations_dir, f"Student_Quiz_Answers_{timestamp}.md")
            
            # Save quiz
            with open(quiz_file, 'w', encoding='utf-8') as f:
                f.write(f"# Student Quiz: {game_info['game_name']}\n\n")
                f.write(f"*Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n")
                f.write(f"**Student Name:** _____________________________\n\n")
                f.write(f"**Date:** _____________________________\n\n")
                f.write("---\n\n")
                f.write(quiz_content)
            
            # Save answers
            with open(answers_file, 'w', encoding='utf-8') as f:
                f.write(f"# Answer Key: {game_info['game_name']}\n\n")
                f.write(f"*Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n")
                f.write("---\n\n")
                f.write(answer_content)
            
            return {
                "quiz": quiz_file,
                "answers": answers_file
            }
        
        except Exception as e:
            print(f"Error creating student quiz: {e}")
            return None
    
    def create_parent_guide(self, game_name):
        """Create a parent guide for the Minecraft Education game."""
        if not self.settings.is_configured():
            return None
        
        game_info = self._gather_all_game_info(game_name)
        
        if not game_info["context"] and not game_info["gameplay"]:
            print("Not enough game information available. Please add context and gameplay first.")
            return None
        
        try:
            from openai import AzureOpenAI
            
            config = self.settings.get_azure_config()
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            objectives_text = "\n".join([f"- {obj}" for obj in game_info["objectives"]]) if game_info["objectives"] else "Not specified"
            
            prompt = f"""Create a comprehensive PARENT GUIDE for the Minecraft Education game "{game_info['game_name']}".

Available Information:

GAME CONTEXT:
{game_info['context']}

GAMEPLAY:
{game_info['gameplay']}

LEARNING OBJECTIVES:
{objectives_text}

LANGUAGE FILE ANALYSIS:
{game_info['lang_analysis'] if game_info['lang_analysis'] else 'Not available'}

ADDITIONAL CONTEXT:
{game_info['document_analysis'] if game_info['document_analysis'] else 'Not available'}

Create a parent guide that includes:

1. OVERVIEW
   - What is this game about?
   - Why is this educational?
   - What will my child learn?

2. HOW TO SUPPORT YOUR CHILD
   - Discussion questions to ask
   - Ways to extend learning at home
   - Connections to real-world topics

3. TECHNICAL INFORMATION
   - How to access the game
   - Required materials or setup
   - Estimated time to complete

4. LEARNING GOALS
   - Clear explanation of objectives
   - Skills being developed
   - How this fits into curriculum

5. FREQUENTLY ASKED QUESTIONS
   - Common concerns addressed
   - Safety and privacy information
   - Who to contact for help

Write in a friendly, accessible tone that helps parents understand the educational value and support their child's learning."""
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "system", "content": "You are an expert in educational communication with parents. Create clear, supportive guides that help parents understand and support their child's learning through Minecraft Education."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.6,
                max_tokens=3000
            )
            
            content = response.choices[0].message.content.strip()
            
            # Save to creations directory
            game_path = os.path.join(self.games_dir, game_name)
            creations_dir = os.path.join(game_path, "creations")
            if not os.path.exists(creations_dir):
                os.makedirs(creations_dir)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = os.path.join(creations_dir, f"Parent_Guide_{timestamp}.md")
            
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"# Parent Guide: {game_info['game_name']}\n\n")
                f.write(f"*Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n")
                f.write("---\n\n")
                f.write(content)
            
            return output_file
        
        except Exception as e:
            print(f"Error creating parent guide: {e}")
            return None
    
    def create_teacher_guide(self, game_name):
        """Create a teacher guide for the Minecraft Education game."""
        if not self.settings.is_configured():
            return None
        
        game_info = self._gather_all_game_info(game_name)
        
        if not game_info["context"] and not game_info["gameplay"]:
            print("Not enough game information available. Please add context and gameplay first.")
            return None
        
        try:
            from openai import AzureOpenAI
            
            config = self.settings.get_azure_config()
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            objectives_text = "\n".join([f"- {obj}" for obj in game_info["objectives"]]) if game_info["objectives"] else "Not specified"
            
            prompt = f"""Create a comprehensive TEACHER GUIDE for the Minecraft Education game "{game_info['game_name']}".

Available Information:

GAME CONTEXT:
{game_info['context']}

GAMEPLAY:
{game_info['gameplay']}

LEARNING OBJECTIVES:
{objectives_text}

LANGUAGE FILE ANALYSIS:
{game_info['lang_analysis'] if game_info['lang_analysis'] else 'Not available'}

ADDITIONAL CONTEXT:
{game_info['document_analysis'] if game_info['document_analysis'] else 'Not available'}

Create a detailed teacher guide that includes:

1. LESSON OVERVIEW
   - Educational goals and standards alignment
   - Grade level and subject areas
   - Time requirements and scheduling suggestions

2. PRE-GAME PREPARATION
   - Technical setup instructions
   - Prerequisite knowledge students need
   - Materials and resources needed
   - Classroom setup recommendations

3. IMPLEMENTATION GUIDE
   - Step-by-step facilitation instructions
   - Key teaching points and interventions
   - Differentiation strategies for diverse learners
   - Pacing recommendations

4. STUDENT SUPPORT STRATEGIES
   - Common challenges and solutions
   - Scaffolding techniques
   - Questions to promote deeper thinking
   - How to support struggling students

5. ASSESSMENT & EVALUATION
   - Formative assessment strategies
   - Observable learning indicators
   - Grading suggestions
   - Extension activities

6. STANDARDS ALIGNMENT
   - Relevant learning standards addressed
   - Cross-curricular connections
   - 21st century skills developed

7. RESOURCES & NEXT STEPS
   - Additional materials
   - Follow-up activities
   - Related lessons

Write in a professional, practical tone that gives teachers actionable guidance."""
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "system", "content": "You are an expert in educational pedagogy and Minecraft Education implementation. Create comprehensive, practical teacher guides that support effective classroom instruction."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.6,
                max_tokens=4000
            )
            
            content = response.choices[0].message.content.strip()
            
            # Save to creations directory
            game_path = os.path.join(self.games_dir, game_name)
            creations_dir = os.path.join(game_path, "creations")
            if not os.path.exists(creations_dir):
                os.makedirs(creations_dir)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = os.path.join(creations_dir, f"Teacher_Guide_{timestamp}.md")
            
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"# Teacher Guide: {game_info['game_name']}\n\n")
                f.write(f"*Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n")
                f.write("---\n\n")
                f.write(content)
            
            return output_file
        
        except Exception as e:
            print(f"Error creating teacher guide: {e}")
            return None
    
    def create_leadership_sheet(self, game_name):
        """Create a school leadership information sheet for the Minecraft Education game."""
        if not self.settings.is_configured():
            return None
        
        game_info = self._gather_all_game_info(game_name)
        
        if not game_info["context"] and not game_info["gameplay"]:
            print("Not enough game information available. Please add context and gameplay first.")
            return None
        
        try:
            from openai import AzureOpenAI
            
            config = self.settings.get_azure_config()
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            objectives_text = "\n".join([f"- {obj}" for obj in game_info["objectives"]]) if game_info["objectives"] else "Not specified"
            
            prompt = f"""Create a concise SCHOOL LEADERSHIP INFORMATION SHEET for the Minecraft Education game "{game_info['game_name']}".

Available Information:

GAME CONTEXT:
{game_info['context']}

GAMEPLAY:
{game_info['gameplay']}

LEARNING OBJECTIVES:
{objectives_text}

LANGUAGE FILE ANALYSIS:
{game_info['lang_analysis'] if game_info['lang_analysis'] else 'Not available'}

ADDITIONAL CONTEXT:
{game_info['document_analysis'] if game_info['document_analysis'] else 'Not available'}

Create a 1-2 page executive summary that includes:

1. PROGRAM OVERVIEW (2-3 paragraphs)
   - What is this educational initiative?
   - Core educational value proposition
   - Target student population

2. STRATEGIC ALIGNMENT
   - How this supports school/district goals
   - Standards and curriculum connections
   - 21st century skills development
   - Equity and access considerations

3. IMPLEMENTATION REQUIREMENTS
   - Technology infrastructure needs
   - Teacher training and support required
   - Time commitment and scheduling
   - Resource allocation

4. EXPECTED OUTCOMES
   - Learning objectives and impacts
   - Student engagement benefits
   - Measurable success indicators
   - Assessment approaches

5. RETURN ON INVESTMENT
   - Educational benefits
   - Student engagement metrics
   - Professional development value
   - Scalability potential

6. RISK MITIGATION
   - Technical support considerations
   - Privacy and safety measures
   - Parent communication strategy

7. RECOMMENDED ACTIONS
   - Next steps for implementation
   - Timeline suggestions
   - Budget considerations
   - Success factors

Write in a concise, executive-level tone focused on strategic value and practical implementation. Use bullet points and clear sections for easy scanning."""
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "system", "content": "You are an expert in educational leadership communication and strategic planning. Create concise, compelling information sheets that help school leaders make informed decisions about educational technology initiatives."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.6,
                max_tokens=3000
            )
            
            content = response.choices[0].message.content.strip()
            
            # Save to creations directory
            game_path = os.path.join(self.games_dir, game_name)
            creations_dir = os.path.join(game_path, "creations")
            if not os.path.exists(creations_dir):
                os.makedirs(creations_dir)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = os.path.join(creations_dir, f"Leadership_Info_Sheet_{timestamp}.md")
            
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"# School Leadership Information Sheet: {game_info['game_name']}\n\n")
                f.write(f"*Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n")
                f.write("---\n\n")
                f.write(content)
            
            return output_file
        
        except Exception as e:
            print(f"Error creating leadership information sheet: {e}")
            return None
    
    def create_curriculum_mapping(self, game_name, country, standards):
        """Create a curriculum standards mapping document."""
        if not self.settings.is_configured():
            return None
        
        game_info = self._gather_all_game_info(game_name)
        
        if not game_info["context"] and not game_info["gameplay"]:
            print("Not enough game information available. Please add context and gameplay first.")
            return None
        
        try:
            from openai import AzureOpenAI
            
            config = self.settings.get_azure_config()
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            objectives_text = "\n".join([f"- {obj}" for obj in game_info["objectives"]]) if game_info["objectives"] else "Not specified"
            standards_text = "\n".join([f"- {std}" for std in standards])
            
            prompt = f"""Create a comprehensive CURRICULUM STANDARDS MAPPING document for the Minecraft Education game "{game_info['game_name']}".

COUNTRY/REGION: {country}

STANDARDS TO MAP:
{standards_text}

GAME INFORMATION:

CONTEXT:
{game_info['context']}

GAMEPLAY:
{game_info['gameplay']}

LEARNING OBJECTIVES:
{objectives_text}

LANGUAGE FILE ANALYSIS:
{game_info['lang_analysis'] if game_info['lang_analysis'] else 'Not available'}

ADDITIONAL CONTEXT:
{game_info['document_analysis'] if game_info['document_analysis'] else 'Not available'}

Create a detailed standards mapping document that includes:

1. EXECUTIVE SUMMARY
   - Overview of standards alignment
   - Key learning domains covered
   - Grade level recommendations

2. DETAILED STANDARDS MAPPING
   For EACH selected standard framework:
   
   a) Standard/Framework Name
   b) Specific standards addressed (with codes/identifiers)
   c) How the game addresses each standard
   d) Evidence of alignment (specific game activities/content)
   e) Depth of coverage (introductory, reinforcement, mastery)
   f) Assessment opportunities

3. CROSS-CURRICULAR CONNECTIONS
   - Multiple subject areas addressed
   - Interdisciplinary learning opportunities
   - 21st century skills developed

4. LEARNING PROGRESSION
   - Prerequisites required
   - Skills developed during gameplay
   - Extensions and follow-up activities

5. DIFFERENTIATION & ACCESSIBILITY
   - How standards can be met at different levels
   - Accommodations for diverse learners
   - Multiple pathways to demonstrate mastery

6. ASSESSMENT ALIGNMENT
   - How to assess standards-based learning
   - Observable evidence of student learning
   - Formative and summative assessment suggestions

7. DOCUMENTATION FOR STAKEHOLDERS
   - Parent communication points
   - Administrative reporting language
   - Evidence for curriculum planning

Be specific and detailed. Include actual standard codes/identifiers where applicable. Provide clear evidence of how game activities align with each standard. Make this document practical for teachers to use for lesson planning and reporting."""
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "system", "content": f"You are an expert in curriculum standards alignment and educational assessment, with deep knowledge of {country} education standards. Create detailed, accurate mappings between learning activities and curriculum standards, using specific standard codes and identifiers where applicable."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.5,
                max_tokens=4000
            )
            
            content = response.choices[0].message.content.strip()
            
            # Save to creations directory
            game_path = os.path.join(self.games_dir, game_name)
            creations_dir = os.path.join(game_path, "creations")
            if not os.path.exists(creations_dir):
                os.makedirs(creations_dir)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = os.path.join(creations_dir, f"Curriculum_Standards_Mapping_{timestamp}.md")
            
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"# Curriculum Standards Mapping: {game_info['game_name']}\n\n")
                f.write(f"**Country/Region:** {country}\n\n")
                f.write(f"**Standards Mapped:**\n")
                for std in standards:
                    f.write(f"- {std}\n")
                f.write(f"\n*Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n")
                f.write("---\n\n")
                f.write(content)
            
            return output_file
        
        except Exception as e:
            print(f"Error creating curriculum mapping: {e}")
            return None
    
    def create_text_complexity_analysis(self, game_name):
        """Create a text complexity analysis with simplification recommendations."""
        if not self.settings.is_configured():
            return None
        
        game_info = self._gather_all_game_info(game_name)
        
        if not game_info["lang_analysis"]:
            print("No language file analysis available. Please extract and analyze language files first.")
            return None
        
        try:
            from openai import AzureOpenAI
            
            config = self.settings.get_azure_config()
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            objectives_text = "\n".join([f"- {obj}" for obj in game_info["objectives"]]) if game_info["objectives"] else "Not specified"
            
            prompt = f"""Analyze the TEXT COMPLEXITY of the in-game language for the Minecraft Education game "{game_info['game_name']}" and provide detailed recommendations for simplification to improve accessibility.

GAME CONTEXT:
{game_info['context'] if game_info['context'] else 'Not available'}

LEARNING OBJECTIVES:
{objectives_text}

LANGUAGE FILE ANALYSIS (NPC Dialogue & In-Game Text):
{game_info['lang_analysis']}

Provide a comprehensive TEXT COMPLEXITY ANALYSIS that includes:

1. EXECUTIVE SUMMARY
   - Overall readability assessment
   - Target age/grade level detected
   - Key complexity issues identified
   - Priority recommendations overview

2. DETAILED COMPLEXITY METRICS
   a) Vocabulary Analysis
      - Academic/technical terms used
      - Domain-specific jargon
      - Word length and frequency
      - Recommended vocabulary level
   
   b) Sentence Structure
      - Average sentence length
      - Complex vs. simple sentences
      - Use of passive voice
      - Conditional and compound structures
   
   c) Reading Level Assessment
      - Estimated grade level (Flesch-Kincaid or similar)
      - Readability score
      - Cognitive load assessment

3. ACCESSIBILITY CONCERNS
   - Barriers for English Language Learners (ELL)
   - Challenges for struggling readers
   - Potential comprehension obstacles
   - Cultural or context-specific references

4. SPECIFIC TEXT EXAMPLES
   For 5-10 representative examples:
   - Original text passage
   - Complexity issues identified
   - Simplified alternative version
   - Explanation of improvements
   - Reading level comparison

5. SIMPLIFICATION STRATEGIES
   a) Vocabulary Recommendations
      - Words to replace with simpler alternatives
      - Technical terms that need definitions
      - Opportunities for visual support
   
   b) Sentence Structure Improvements
      - Breaking long sentences
      - Converting passive to active voice
      - Simplifying complex grammar
   
   c) Content Organization
      - Adding headers or structure
      - Breaking information into chunks
      - Using bullet points or lists
      - Adding context or scaffolding

6. DIFFERENTIATION SUGGESTIONS
   - Support for multiple reading levels
   - Options for text-to-speech integration
   - Visual aids or glossaries
   - Scaffolded text versions

7. IMPLEMENTATION GUIDE
   - Priority order for changes
   - Quick wins vs. major revisions
   - Testing recommendations
   - Quality assurance checklist

8. BEFORE & AFTER EXAMPLES
   Provide 3-5 complete examples showing:
   - Original complex text
   - Revised simplified version
   - Reading level improvement
   - Maintained educational value

Be specific and actionable. Provide actual text examples from the game. Focus on maintaining educational integrity while improving accessibility. Consider the target audience and learning objectives when making recommendations."""
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "system", "content": "You are an expert in educational linguistics, readability analysis, and accessible content design. You specialize in analyzing text complexity for educational games and providing practical recommendations to improve accessibility while maintaining learning objectives. Use readability formulas, cognitive load theory, and UDL principles in your analysis."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.5,
                max_tokens=4000
            )
            
            content = response.choices[0].message.content.strip()
            
            # Save to creations directory
            game_path = os.path.join(self.games_dir, game_name)
            creations_dir = os.path.join(game_path, "creations")
            if not os.path.exists(creations_dir):
                os.makedirs(creations_dir)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = os.path.join(creations_dir, f"Text_Complexity_Analysis_{timestamp}.md")
            
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(f"# Text Complexity Analysis: {game_info['game_name']}\n\n")
                f.write(f"*Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*\n\n")
                f.write("---\n\n")
                f.write("## Purpose\n\n")
                f.write("This analysis evaluates the complexity of in-game text and provides recommendations ")
                f.write("to simplify language for improved accessibility while maintaining educational value.\n\n")
                f.write("---\n\n")
                f.write(content)
            
            return output_file
        
        except Exception as e:
            print(f"Error creating text complexity analysis: {e}")
            return None
    
    def analyze_lang_file_with_ai(self, game_name):
        """Analyze the extracted language file using Azure OpenAI and save in standardized format."""
        if not self.settings.is_configured():
            print("Azure OpenAI is not configured!")
            return None
        
        game_path = os.path.join(self.games_dir, game_name)
        lang_dir = os.path.join(game_path, "lang")
        
        # Find the JSON version of the lang file
        if not os.path.exists(lang_dir):
            return None
            
        json_files = [f for f in os.listdir(lang_dir) if f.endswith('.json') and f != 'extraction_analysis.json']
        
        if not json_files:
            return None
        
        json_file = os.path.join(lang_dir, json_files[0])
        
        try:
            with open(json_file, 'r', encoding='utf-8') as f:
                lang_data = json.load(f)
            
            # Create a sample of the lang data for analysis
            sample_entries = dict(list(lang_data.items())[:50])  # First 50 entries
            sample_text = json.dumps(sample_entries, indent=2)
            
            from openai import AzureOpenAI
            
            config = self.settings.get_azure_config()
            
            client = AzureOpenAI(
                api_key=config['api_key'],
                api_version=config['api_version'],
                azure_endpoint=config['endpoint']
            )
            
            prompt = f"""Analyze this Minecraft Education language file data and provide a structured analysis in the following format:

1. NARRATIVE ELEMENTS:
   - Main story or theme
   - Key characters or NPCs
   - Plot points or story progression

2. EDUCATIONAL CONTENT:
   - Subject areas covered
   - Key concepts and terms
   - Learning activities mentioned

3. GAME OBJECTIVES:
   - Primary goals
   - Quest structure
   - Success criteria

4. INSTRUCTIONAL TEXT:
   - Player guidance
   - Tutorial elements
   - Help text

5. EDUCATIONAL FOCUS:
   - Grade level appropriateness
   - Curriculum alignment
   - Overall educational theme

Language file sample (showing {len(sample_entries)} of {len(lang_data)} entries):
{sample_text}"""
            
            response = client.chat.completions.create(
                model=config['deployment'],
                messages=[
                    {"role": "system", "content": "You are an educational game content analyst specializing in Minecraft Education. Provide structured, detailed analysis of language files to extract educational context and gameplay information."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.5,
                max_tokens=2000
            )
            
            analysis_text = response.choices[0].message.content.strip()
            
            # Save analysis in standardized format
            analysis_data = {
                "analyzed_at": datetime.now().isoformat(),
                "total_entries": len(lang_data),
                "sample_size": len(sample_entries),
                "analysis": analysis_text,
                "summary": self._extract_summary(analysis_text)
            }
            
            # Save to game info folder
            self.save_game_info(game_name, "lang_analysis", analysis_data)
            
            # Update metadata
            self._update_lang_metadata(game_name, analyzed=True)
            
            return analysis_text
        
        except Exception as e:
            print(f"Error analyzing lang file: {e}")
            return None
    
    def _extract_summary(self, analysis_text):
        """Extract a brief summary from the analysis text."""
        lines = analysis_text.split('\n')
        summary_lines = []
        
        for line in lines[:15]:  # First 15 lines as summary
            if line.strip() and not line.strip().startswith('#'):
                summary_lines.append(line.strip())
                if len(summary_lines) >= 5:
                    break
        
        return ' '.join(summary_lines)[:300] + '...' if summary_lines else "Analysis available"
    
    def _update_lang_metadata(self, game_name, analyzed=False):
        """Update language file metadata with analysis status."""
        game_path = os.path.join(self.games_dir, game_name)
        metadata_file = os.path.join(game_path, "metadata.json")
        
        try:
            with open(metadata_file, 'r') as f:
                metadata = json.load(f)
            
            if "lang_file" in metadata:
                metadata["lang_file"]["ai_analyzed"] = analyzed
                metadata["lang_file"]["analyzed_at"] = datetime.now().isoformat()
            
            metadata["modified"] = datetime.now().isoformat()
            
            with open(metadata_file, 'w') as f:
                json.dump(metadata, f, indent=4)
        except (json.JSONDecodeError, IOError) as e:
            print(f"Error updating metadata: {e}")
